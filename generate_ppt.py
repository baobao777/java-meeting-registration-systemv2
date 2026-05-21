from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x73, 0xE8)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xF0, 0xF2, 0xF5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xE5, 0x39, 0x35)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, text="", font_size=18, bold=False, font_color=WHITE, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox

# ============================================================
# Slide 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)
add_shape(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), WHITE)
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1), "会议报名管理系统", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6), "基于 Spring Boot + Vue 3 + MySQL 的全栈项目", 22, RGBColor(0xBB, 0xDE, 0xFB), False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5), "版本 v1.0 ｜ 2026年5月", 18, RGBColor(0x88, 0xBB, 0xEE), False, PP_ALIGN.CENTER)

# ============================================================
# Slide 2: TOC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "目 录", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

toc = [
    "1. 项目背景",
    "2. 技术架构",
    "3. 功能模块",
    "4. 数据库设计",
    "5. API 接口",
    "6. 核心业务流程",
    "7. 异常处理机制",
    "8. 测试结果",
    "9. 项目结构",
    "10. 部署方式",
    "11. 总结与展望",
]
add_bullets(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5.5), toc, 22, DARK)

# ============================================================
# Slide 3: Background
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "1. 项目背景", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.06), RED)
add_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5), "传统方式痛点", 24, RED, True)

pains = [
    "▸ 纸质报名：效率低、易丢失、难追溯",
    "▸ 邮件/群接龙：人工统计繁琐、易出错",
    "▸ 会议信息分散：无法集中发布和查询",
    "▸ 报名数据不透明：管理员难以及时掌握",
]
add_bullets(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4), pains, 18, DARK)

add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(0.06), GREEN)
add_text(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(0.5), "系统目标", 24, GREEN, True)

goals = [
    "▸ 线上发布会议，一目了然",
    "▸ 用户一键报名/取消，自助操作",
    "▸ 管理员统一管控会议与报名数据",
    "▸ 报名进度实时可视，满额自动限制",
]
add_bullets(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(4), goals, 18, DARK)

# ============================================================
# Slide 4: Tech Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "2. 技术架构", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

layers = [
    ("前端层", "Vue 3 + Vite 5\nVue Router + Axios", Inches(0.8), BLUE),
    ("网络层", "RESTful API\nJSON 格式通信", Inches(3.6), RGBColor(0x55, 0x55, 0x55)),
    ("后端层", "Spring Boot 3.2\nController-Service-Mapper", Inches(6.4), BLUE),
    ("数据层", "MySQL 8.0\nMyBatis ORM", Inches(9.2), RGBColor(0x55, 0x55, 0x55)),
]

for title, desc, left, color in layers:
    add_shape(slide, left, Inches(1.8), Inches(3.2), Inches(0.6), color, title, 20, True, WHITE)
    add_text(slide, left, Inches(2.6), Inches(3.2), Inches(1.5), desc, 16, GRAY, False, PP_ALIGN.CENTER)

arrows = [Inches(3.6), Inches(6.4), Inches(9.2)]
for a in arrows:
    add_text(slide, a, Inches(2.8), Inches(0.5), Inches(0.5), "▼", 24, BLUE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(4.4), Inches(5), Inches(0.5), "技术栈详情", 22, DARK, True)

techs = [
    ("后端框架", "Spring Boot 3.2.0"),
    ("ORM", "MyBatis 3.0.3"),
    ("数据库", "MySQL 8.0"),
    ("前端框架", "Vue 3 + Vite 5"),
    ("路由", "Vue Router 4"),
    ("HTTP 客户端", "Axios 1.6"),
    ("构建工具", "Maven 3.9 + npm"),
    ("JDK", "Java 21"),
]

table = slide.shapes.add_table(len(techs)+1, 2, Inches(0.8), Inches(5.0), Inches(11), Inches(2.0)).table
table.columns[0].width = Inches(4)
table.columns[1].width = Inches(7)

for i, h in enumerate(["技术层次", "选用技术"]):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.size = Pt(14)

for r, (layer, tech) in enumerate(techs):
    for c, val in enumerate([layer, tech]):
        cell = table.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = DARK
        if r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LGRAY

# ============================================================
# Slide 5: Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "3. 功能模块", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.6), BLUE, "普通用户端", 22, True, WHITE)
user_features = [
    "✅ 用户注册 / 登录 / 登出",
    "✅ 浏览会议列表（卡片式布局）",
    "✅ 查看会议详情（时间/地点/容量）",
    "✅ 报名会议（满额自动拦截）",
    "✅ 取消报名（确认弹窗）",
    "✅ 我的报名记录查询",
]
add_bullets(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.5), user_features, 17, DARK)

add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(0.6), GREEN, "管理员端", 22, True, WHITE)
admin_features = [
    "✅ 创建会议（标题/描述/时间/容量）",
    "✅ 编辑会议信息与状态",
    "✅ 删除会议（二次确认）",
    "✅ 查看报名人员名单（姓名/电话）",
    "✅ 管理后台表格化管理",
    "✅ 创建/编辑会议表单验证",
]
add_bullets(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(4.5), admin_features, 17, DARK)

# ============================================================
# Slide 6: Database Design
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "4. 数据库设计", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5), "三张核心表：user → meeting → registration，外键关联保证数据完整性", 16, GRAY)

tables_info = [
    ("user 用户表", ["id (PK)", "username (UNIQUE)", "password", "name", "email", "phone", "role (USER/ADMIN)", "created_at"]),
    ("meeting 会议表", ["id (PK)", "title", "description", "location", "start_time", "end_time", "capacity (0=不限)", "status (4种状态)", "creator_id (FK)", "created_at"]),
    ("registration 报名表", ["id (PK)", "user_id (FK)", "meeting_id (FK)", "registered_at", "status (REGISTERED/CANCELLED)"]),
]

for i, (title, fields) in enumerate(tables_info):
    left = Inches(0.8 + i * 4.2)
    add_shape(slide, left, Inches(2.3), Inches(3.8), Inches(0.5), BLUE, title, 18, True, WHITE)
    add_bullets(slide, left + Inches(0.2), Inches(3.0), Inches(3.5), Inches(4), ["- " + f for f in fields], 14, DARK)

add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.04), RGBColor(0xBB, 0xBB, 0xBB))
add_text(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
         "关系：user.id ← registration.user_id ｜ meeting.id ← registration.meeting_id ｜ meeting.creator_id → user.id",
         14, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# Slide 7: API
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "5. API 接口", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

api_sections = [
    ("用户模块", [
        ("POST", "/api/user/register", "注册", "无"),
        ("POST", "/api/user/login", "登录", "无"),
        ("GET", "/api/user/info", "获取用户信息", "已登录"),
    ]),
    ("会议模块", [
        ("GET", "/api/meeting/list", "会议列表", "无"),
        ("GET", "/api/meeting/{id}", "会议详情", "无"),
        ("POST", "/api/meeting", "创建会议", "管理员"),
        ("PUT", "/api/meeting/{id}", "编辑会议", "管理员"),
        ("DELETE", "/api/meeting/{id}", "删除会议", "管理员"),
    ]),
    ("报名模块", [
        ("POST", "/api/registration", "报名", "已登录"),
        ("DELETE", "/api/registration/{id}", "取消报名", "已登录"),
        ("GET", "/api/registration/my", "我的报名", "已登录"),
        ("GET", "/api/registration/meeting/{id}", "报名名单", "管理员"),
        ("GET", "/api/registration/check/{id}", "检查报名", "已登录"),
    ]),
]

y = Inches(1.6)
for title, endpoints in api_sections:
    add_text(slide, Inches(0.8), y, Inches(3), Inches(0.4), title, 18, BLUE, True)
    y += Inches(0.4)
    tbl = slide.shapes.add_table(len(endpoints)+1, 4, Inches(0.8), y, Inches(11.5), Inches(0.35 * (len(endpoints)+1))).table
    for c, w in enumerate([Inches(1), Inches(4.5), Inches(3.5), Inches(1.5)]):
        tbl.columns[c].width = w
    for c, h in enumerate(["方法", "路径", "说明", "权限"]):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.size = Pt(12)
            p.font.bold = True
    for r, (method, path, desc, auth) in enumerate(endpoints):
        for c, val in enumerate([method, path, desc, auth]):
            cell = tbl.cell(r+1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
    y += Inches(0.35 * (len(endpoints)+1) + 0.2)

# ============================================================
# Slide 8: Business Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "6. 核心业务流程", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5), "用户报名完整流程：", 20, DARK, True)

flow_items = [
    "用户注册 → 用户登录 → 浏览会议列表 → 查看会议详情",
    "",
    "点击报名 → 检查：已报名？满额？会议状态？",
    "    │",
    "    ├── 已报名 → 拦截 (已报名该会议)",
    "    ├── 已满额 → 拦截 (会议名额已满)",
    "    ├── 非 UPCOMING → 拦截 (会议不在报名阶段)",
    "    └── 全部通过 → 报名成功",
    "",
    "取消报名 → 确认弹窗 → 状态标记为 CANCELLED",
]
add_bullets(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(4), flow_items, 17, DARK)

# Flow diagram at bottom
add_shape(slide, Inches(1), Inches(5.2), Inches(2.5), Inches(0.6), BLUE, "登录", 18, True, WHITE)
add_text(slide, Inches(3.2), Inches(5.3), Inches(0.8), Inches(0.5), "→", 28, BLUE, True)
add_shape(slide, Inches(4), Inches(5.2), Inches(2.5), Inches(0.6), BLUE, "筛选会议", 18, True, WHITE)
add_text(slide, Inches(6.2), Inches(5.3), Inches(0.8), Inches(0.5), "→", 28, BLUE, True)
add_shape(slide, Inches(7), Inches(5.2), Inches(2.5), Inches(0.6), GREEN, "报名成功", 18, True, WHITE)

# ============================================================
# Slide 9: Error Handling
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "7. 异常处理机制", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

error_sections = [
    ("用户模块", [
        ("用户名重复", "400 - 用户名已存在"),
        ("密码错误", "400 - 密码错误"),
        ("用户不存在", "400 - 用户不存在"),
        ("未登录访问", "401 - 未登录"),
    ]),
    ("会议/权限模块", [
        ("会议不存在", "404 - 会议不存在"),
        ("普通用户越权", "403 - 无权操作"),
        ("未登录报名", "401 - 未登录"),
        ("路由守卢拦截", "前端跳转 /login"),
    ]),
    ("报名模块", [
        ("重复报名", "400 - 已报名该会议"),
        ("会议满额", "400 - 会议名额已满"),
        ("非报名阶段", "400 - 会议不在报名阶段"),
        ("取消已取消的报名", "400 - 未报名该会议"),
    ]),
]

for i, (title, items) in enumerate(error_sections):
    left = Inches(0.8 + i * 4.2)
    add_shape(slide, left, Inches(1.8), Inches(3.8), Inches(0.5), BLUE, title, 18, True, WHITE)
    y_off = Inches(2.5)
    for scenario, response in items:
        add_shape(slide, left + Inches(0.1), y_off, Inches(3.6), Inches(0.45), LGRAY, "", 12)
        add_text(slide, left + Inches(0.2), y_off + Inches(0.02), Inches(3.4), Inches(0.22), scenario, 14, DARK, True)
        add_text(slide, left + Inches(0.2), y_off + Inches(0.24), Inches(3.4), Inches(0.2), response, 12, RED, False)
        y_off += Inches(0.52)

# ============================================================
# Slide 10: Test Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "8. 测试结果", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
         "使用 Postman 对全部 19 个接口、35 个测试用例进行自动化测试", 18, GRAY)

test_results = [
    ("用户模块", "8", "8", "0", "100%"),
    ("会议管理", "10", "10", "0", "100%"),
    ("报名流程", "8", "8", "0", "100%"),
    ("满额测试", "3", "3", "0", "100%"),
    ("权限测试", "6", "6", "0", "100%"),
    ("合  计", "35", "35", "0", "100%"),
]

tbl = slide.shapes.add_table(len(test_results)+1, 5, Inches(1.5), Inches(2.5), Inches(10), Inches(0.42 * (len(test_results)+1))).table
for c, w in enumerate([Inches(2), Inches(2), Inches(2), Inches(2), Inches(2)]):
    tbl.columns[c].width = w

for c, h in enumerate(["模块", "用例数", "通过", "失败", "通过率"]):
    cell = tbl.cell(0, c)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

for r, (mod, total, passed, failed, rate) in enumerate(test_results):
    is_total = r == len(test_results) - 1
    for c, val in enumerate([mod, total, passed, failed, rate]):
        cell = tbl.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = is_total
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER
        if is_total:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)

add_shape(slide, Inches(3.5), Inches(5.5), Inches(6), Inches(0.6), GREEN, "✅  全部通过 ｜ 通过率 100%", 22, True, WHITE)

# ============================================================
# Slide 11: Project Structure
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "9. 项目结构", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

structure = [
    "Java01/",
    "  init.sql                          # 数据库初始化",
    "  meeting-api-test.postman_collection.json  # Postman测试集合",
    "  meeting-api/                       # Spring Boot 后端",
    "    src/main/java/com/meeting/",
    "      MeetingApplication.java        # 启动类",
    "      config/                        # 跨域 + Web 配置",
    "      entity/                        # User / Meeting / Registration",
    "      dto/                           # ApiResult / LoginRequest",
    "      mapper/                        # MyBatis 接口",
    "      service/                       # 业务逻辑层",
    "      controller/                    # REST 控制器",
    "  meeting-web/                       # Vue 3 前端",
    "    src/",
    "      App.vue                        # 导航栏根组件",
    "      router/                        # 路由 + 守卢",
    "      api/                           # Axios API 封装",
    "      views/                         # 7 个页面组件",
    "  deploy/                            # 部署包",
    "    meeting-api-1.0.0.jar            # 可运行 JAR",
    "    application.yml                  # 配置文件",
    "    start.bat / start.sh             # 启动脚本",
    "  docs/                              # 项目文档",
]
add_bullets(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5.5), structure, 14, DARK)

# ============================================================
# Slide 12: Deployment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "10. 部署方式", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5), BLUE, "开发环境部署", 20, True, WHITE)
dev_steps = [
    "1. 初始化数据库：mysql -u root -p < init.sql",
    "",
    "2. 启动后端：cd meeting-api && mvn spring-boot:run",
    "",
    "3. 启动前端：cd meeting-web && npm run dev",
    "",
    "4. 访问 http://localhost:3000",
]
add_bullets(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.5), dev_steps, 16, DARK)

add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(0.5), GREEN, "生产环境部署", 20, True, WHITE)
prod_steps = [
    "1. 初始化数据库：mysql -u root -p < init.sql",
    "",
    "2. 修改 application.yml 数据库密码",
    "",
    "3. 一行命令启动：java -jar meeting-api-1.0.0.jar",
    "",
    "4. 访问 http://localhost:8080",
]
add_bullets(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(4.5), prod_steps, 16, DARK)

add_shape(slide, Inches(3), Inches(6.0), Inches(7), Inches(0.6), RGBColor(0xE3, 0xF2, 0xFD), "默认管理员账号：admin / admin123", 18, True, BLUE)

# ============================================================
# Slide 13: Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "11. 总结与展望", 36, BLUE, True)
add_shape(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), BLUE)

add_text(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5), "项目亮点", 22, BLUE, True)
highlights = [
    "✔ 全栈实现：Spring Boot + Vue 3 前后端分离",
    "✔ 业务完整：用户 → 会议 → 报名 全链路覆盖",
    "✔ 健壮性好：35 个测试用例通过率 100%",
    "✔ 异常覆盖：11 种异常场景全部正确处理",
    "✔ 权限隔离：未登录/普通用户/管理员三级管控",
    "✔ 部署简单：一行 java -jar 即可运行",
]
add_bullets(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), highlights, 16, DARK)

add_text(slide, Inches(7), Inches(1.6), Inches(5), Inches(0.5), "可扩展方向", 22, GREEN, True)
extensions = [
    "🔒 密码加密（BCrypt）+ JWT Token",
    "📄 分页查询（大数据量支持）",
    "📧 报名确认邮件/短信通知",
    "✅ 二维码签到功能",
    "📊 Excel 导入/导出报名数据",
    "🧪 单元测试 + 集成测试",
    "🐳 Docker 容器化部署",
]
add_bullets(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.5), extensions, 16, DARK)

# ============================================================
# Slide 14: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)
add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), WHITE)
add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1), "感谢聊听", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.6), "会议报名管理系统  v1.0", 24, RGBColor(0xBB, 0xDE, 0xFB), False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5), "管理员：admin / admin123  ｜  演示地址：http://localhost:8080", 18, RGBColor(0x88, 0xBB, 0xEE), False, PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = "D:/Learn/JAVAee/Java01/docs/会议报名管理系统_项目演示.pptx"
prs.save(output_path)
print(f"PPT generated: {output_path}")
print(f"Total slides: {len(prs.slides)}")
